"""Report generation — PDF and PowerPoint, built from the dashboard's current filter.

The browser posts the commits it is *actually showing* plus the criteria that
produced them. Every number in the report is then derived here from that one
list, so the report and the screen can never disagree — there is no second
implementation of the filters to drift out of sync.
"""

from __future__ import annotations

from collections import Counter, defaultdict
from dataclasses import dataclass, field
from datetime import datetime
from io import BytesIO
from typing import Any, Iterable

# Match the dashboard's palette so the report reads as the same product.
ACCENT = "#0a7568"
ACCENT_DARK = "#085f55"
ACCENT_SOFT = "#e3f5f2"
INK = "#17191c"
MUTED = "#5c636b"
BAND = "#eff3f7"
LINE = "#d9dfe6"

# Beyond this the detail table stops being readable; we cut it and say so
# rather than silently truncating.
MAX_DETAIL_ROWS = 400


# --------------------------------------------------------------------------
# aggregation
# --------------------------------------------------------------------------


@dataclass
class Rollup:
    """Everything both output formats need, derived once from the commit list."""

    commits: list[dict[str, Any]]
    criteria: dict[str, Any]
    summary: dict[str, int] = field(default_factory=dict)
    by_repo: list[tuple[str, dict[str, Any]]] = field(default_factory=list)
    by_service: list[tuple[str, dict[str, Any]]] = field(default_factory=list)
    by_author: list[tuple[str, dict[str, Any]]] = field(default_factory=list)
    by_day: list[tuple[str, int]] = field(default_factory=list)
    truncated: int = 0


def _service_keys(commit: dict[str, Any]) -> list[str]:
    """(repo, folder) pairs — a `backend` in two repos is two services."""
    repo = commit.get("repo") or "?"
    folders = commit.get("folders") or []
    return [f"{repo} / {f}" for f in folders]


def _blank_bucket() -> dict[str, Any]:
    return {"commits": 0, "files": 0, "authors": set(), "services": set(), "repos": set()}


def aggregate(commits: list[dict[str, Any]], criteria: dict[str, Any]) -> Rollup:
    roll = Rollup(commits=commits, criteria=criteria)

    repos: dict[str, dict[str, Any]] = defaultdict(_blank_bucket)
    services: dict[str, dict[str, Any]] = defaultdict(_blank_bucket)
    authors: dict[str, dict[str, Any]] = defaultdict(_blank_bucket)
    days: Counter = Counter()

    all_services: set[str] = set()
    all_branches: set[str] = set()

    for c in commits:
        repo = c.get("repo") or "?"
        author = c.get("author_name") or "unknown"
        files = int(c.get("files_changed") or 0)
        svc = _service_keys(c)
        day = (c.get("date") or "")[:10]

        all_services.update(svc)
        all_branches.update(f"{repo}#{b}" for b in (c.get("branches") or []))
        if day:
            days[day] += 1

        r = repos[repo]
        r["commits"] += 1
        r["files"] += files
        r["authors"].add(author)
        r["services"].update(svc)

        a = authors[author]
        a["commits"] += 1
        a["files"] += files
        a["repos"].add(repo)
        a["services"].update(svc)

        for s in svc:
            b = services[s]
            b["commits"] += 1
            b["files"] += files
            b["authors"].add(author)
            b["repos"].add(repo)

    roll.summary = {
        "commits": len(commits),
        "repositories": len(repos),
        "services": len(all_services),
        "branches": len(all_branches),
        "contributors": len(authors),
        "off_default": sum(1 for c in commits if not c.get("on_default")),
        "files_changed": sum(int(c.get("files_changed") or 0) for c in commits),
    }

    by_commits = lambda kv: (-kv[1]["commits"], kv[0])  # noqa: E731
    roll.by_repo = sorted(repos.items(), key=by_commits)
    roll.by_service = sorted(services.items(), key=by_commits)
    roll.by_author = sorted(authors.items(), key=by_commits)
    roll.by_day = sorted(days.items(), reverse=True)

    if len(commits) > MAX_DETAIL_ROWS:
        roll.truncated = len(commits) - MAX_DETAIL_ROWS

    return roll


def _criteria_rows(criteria: dict[str, Any]) -> list[tuple[str, str]]:
    """Only the filters that are actually set, in drill-down order."""
    order = [
        ("comparison", "Comparison"),
        ("merge_base", "Merge base"),
        ("range", "Date range"),
        ("source", "Source"),
        ("repository", "Repository"),
        ("service", "Service / folder"),
        ("branch", "Branch"),
        ("author", "Author"),
        ("show", "Showing"),
        # Distinct files in the whole diff. The Summary's "Files changed" counts
        # per-commit touches instead, so a file edited twice counts twice —
        # different measures, hence different labels.
        ("files", "Net diff vs base"),
        ("search", "Search"),
        ("grouped_by", "Grouped by"),
    ]
    rows = []
    for key, label in order:
        value = criteria.get(key)
        if value:
            rows.append((label, str(value)))
    return rows


def _fmt_date(iso: str | None) -> str:
    if not iso:
        return ""
    try:
        return datetime.fromisoformat(iso.replace("Z", "+00:00")).strftime("%d %b %Y")
    except ValueError:
        return str(iso)[:10]


def _fmt_datetime(iso: str | None) -> str:
    if not iso:
        return ""
    try:
        return datetime.fromisoformat(iso.replace("Z", "+00:00")).strftime("%d %b %Y, %H:%M UTC")
    except ValueError:
        return str(iso)


def _clip(text: str, limit: int) -> str:
    text = " ".join(str(text or "").split())
    return text if len(text) <= limit else text[: limit - 1].rstrip() + "…"


def _join(values: Iterable[str], limit: int = 3) -> str:
    vals = list(values)
    if not vals:
        return "—"
    if len(vals) <= limit:
        return ", ".join(vals)
    return ", ".join(vals[:limit]) + f" +{len(vals) - limit}"


def report_filename(criteria: dict[str, Any], extension: str) -> str:
    since = (criteria.get("since") or "")[:10]
    until = (criteria.get("until") or "")[:10] or "now"
    scope = criteria.get("repository") or criteria.get("source") or "all-repos"
    scope = "".join(ch if ch.isalnum() or ch in "-_" else "-" for ch in str(scope)).strip("-")
    return f"repo-changes_{scope}_{since}_to_{until}.{extension}"


# --------------------------------------------------------------------------
# PDF
# --------------------------------------------------------------------------


def build_pdf(roll: Rollup) -> bytes:
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_LEFT
    from reportlab.lib.pagesizes import A4, landscape
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import mm
    from reportlab.platypus import (
        KeepTogether,
        PageBreak,
        Paragraph,
        SimpleDocTemplate,
        Spacer,
        Table,
        TableStyle,
    )
    from xml.sax.saxutils import escape

    buf = BytesIO()
    page = landscape(A4)
    doc = SimpleDocTemplate(
        buf,
        pagesize=page,
        leftMargin=14 * mm,
        rightMargin=14 * mm,
        topMargin=14 * mm,
        bottomMargin=16 * mm,
        title="Repository Change Report",
        author="Repo Change Dashboard",
    )
    width = page[0] - doc.leftMargin - doc.rightMargin

    base = getSampleStyleSheet()
    st = {
        "h1": ParagraphStyle("h1", parent=base["Title"], fontSize=22, leading=26,
                             textColor=colors.HexColor(INK), alignment=TA_LEFT, spaceAfter=2),
        "sub": ParagraphStyle("sub", parent=base["Normal"], fontSize=11, leading=15,
                              textColor=colors.HexColor(MUTED), spaceAfter=10),
        "h2": ParagraphStyle("h2", parent=base["Heading2"], fontSize=12.5, leading=16,
                             textColor=colors.HexColor(ACCENT_DARK), spaceBefore=12, spaceAfter=5),
        "cell": ParagraphStyle("cell", parent=base["Normal"], fontSize=8, leading=10.5,
                               textColor=colors.HexColor(INK)),
        "note": ParagraphStyle("note", parent=base["Normal"], fontSize=8.5, leading=12,
                               textColor=colors.HexColor(MUTED), spaceBefore=4),
    }

    def table(data: list[list[Any]], widths: list[float], align_right: list[int] = ()) -> Table:
        t = Table(data, colWidths=widths, repeatRows=1, hAlign="LEFT")
        style = [
            # Numeric headers follow their column's alignment, otherwise the
            # label floats left while the figures sit right.
            *[("ALIGN", (c, 0), (c, 0), "RIGHT") for c in align_right],
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(ACCENT_SOFT)),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.HexColor(ACCENT_DARK)),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE", (0, 0), (-1, -1), 8),
            ("LEADING", (0, 0), (-1, -1), 10.5),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("TOPPADDING", (0, 0), (-1, -1), 4),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
            ("LEFTPADDING", (0, 0), (-1, -1), 6),
            ("RIGHTPADDING", (0, 0), (-1, -1), 6),
            ("LINEBELOW", (0, 0), (-1, -2), 0.4, colors.HexColor(LINE)),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor(BAND)]),
            ("BOX", (0, 0), (-1, -1), 0.6, colors.HexColor(LINE)),
        ]
        for col in align_right:
            style.append(("ALIGN", (col, 1), (col, -1), "RIGHT"))
        t.setStyle(TableStyle(style))
        return t

    def p(text: str) -> Paragraph:
        return Paragraph(escape(str(text or "")), st["cell"])

    story: list[Any] = []
    crit = roll.criteria

    # --- cover -----------------------------------------------------------
    story.append(Paragraph("Repository Change Report", st["h1"]))
    span = f"{_fmt_date(crit.get('since'))} – {_fmt_date(crit.get('until')) or 'today'}"
    story.append(Paragraph(f"Commit activity, {span}", st["sub"]))

    rows = _criteria_rows(crit)
    if rows:
        story.append(Paragraph("Filter criteria", st["h2"]))
        story.append(
            table(
                [["Field", "Value"]] + [[p(k), p(v)] for k, v in rows],
                [45 * mm, width - 45 * mm],
            )
        )

    # --- summary ---------------------------------------------------------
    s = roll.summary
    story.append(Paragraph("Summary", st["h2"]))
    story.append(
        table(
            [
                ["Commits", "Repositories", "Services", "Branches", "Contributors",
                 "Not on default", "File changes"],
                [str(s["commits"]), str(s["repositories"]), str(s["services"]),
                 str(s["branches"]), str(s["contributors"]), str(s["off_default"]),
                 str(s["files_changed"])],
            ],
            [width / 7] * 7,
        )
    )

    # --- breakdowns ------------------------------------------------------
    if roll.by_repo:
        story.append(Paragraph("Activity by repository", st["h2"]))
        story.append(
            table(
                [["Repository", "Commits", "Services touched", "Contributors", "Files"]]
                + [
                    [p(name), str(v["commits"]), p(_join(sorted(x.split(" / ")[-1] for x in v["services"]), 6)),
                     str(len(v["authors"])), str(v["files"])]
                    for name, v in roll.by_repo
                ],
                [width * 0.28, width * 0.08, width * 0.40, width * 0.11, width * 0.13],
                align_right=[1, 3, 4],
            )
        )

    if roll.by_service:
        story.append(Paragraph("Activity by service / folder", st["h2"]))
        story.append(
            table(
                [["Service (repository / folder)", "Commits", "Contributors", "Files"]]
                + [
                    [p(name), str(v["commits"]), str(len(v["authors"])), str(v["files"])]
                    for name, v in roll.by_service
                ],
                [width * 0.55, width * 0.15, width * 0.15, width * 0.15],
                align_right=[1, 2, 3],
            )
        )

    if roll.by_author:
        story.append(Paragraph("Activity by contributor", st["h2"]))
        story.append(
            table(
                [["Contributor", "Commits", "Repositories", "Services", "Files"]]
                + [
                    [p(name), str(v["commits"]), str(len(v["repos"])), str(len(v["services"])), str(v["files"])]
                    for name, v in roll.by_author
                ],
                [width * 0.40, width * 0.15, width * 0.15, width * 0.15, width * 0.15],
                align_right=[1, 2, 3, 4],
            )
        )

    if roll.by_day:
        story.append(Paragraph("Activity by day", st["h2"]))
        cols = 6
        cells = [f"{_fmt_date(d)}  ({n})" for d, n in roll.by_day]
        grid = [cells[i : i + cols] for i in range(0, len(cells), cols)]
        grid = [row + [""] * (cols - len(row)) for row in grid]
        story.append(
            table([["Date (commits)"] + [""] * (cols - 1)] + [[p(c) for c in row] for row in grid],
                  [width / cols] * cols)
        )

    # --- detail ----------------------------------------------------------
    story.append(PageBreak())
    story.append(Paragraph("Commit detail", st["h2"]))

    detail = roll.commits[:MAX_DETAIL_ROWS]
    head = ["Date", "Repository", "Service", "Branch", "Author", "SHA", "Files", "Message"]
    body = [
        [
            # ISO dates here rather than "04 Aug 2026": half the width, unambiguous,
            # and they sort correctly if the table is pasted into a spreadsheet.
            p((c.get("date") or "")[:10]),
            p(_clip(c.get("repo"), 26)),
            p(_join(c.get("folders") or [], 2)),
            p(_join(c.get("branches") or [], 2)),
            p(_clip(c.get("author_name"), 18)),
            p((c.get("sha") or "")[:7]),
            str(c.get("files_changed") or 0),
            p(_clip(c.get("title"), 90)),
        ]
        for c in detail
    ]
    story.append(
        table(
            [head] + body,
            # Sized so Date / SHA / Files / Author never wrap; Message takes the
            # slack and is the only column allowed to run to a second line.
            [
                width * 0.075, width * 0.160, width * 0.140, width * 0.100,
                width * 0.130, width * 0.060, width * 0.040, width * 0.295,
            ],
            align_right=[6],
        )
    )
    if roll.truncated:
        story.append(
            Paragraph(
                f"Showing the {MAX_DETAIL_ROWS} most recent commits of {len(roll.commits)}. "
                f"{roll.truncated} older commits are counted in every total above but not "
                "listed individually — narrow the date range or the filters to list them all.",
                st["note"],
            )
        )

    generated = _fmt_datetime(crit.get("generated_at"))

    def furniture(canvas, doc_):
        canvas.saveState()
        canvas.setStrokeColor(colors.HexColor(LINE))
        canvas.setLineWidth(0.5)
        canvas.line(doc_.leftMargin, 11 * mm, page[0] - doc_.rightMargin, 11 * mm)
        canvas.setFont("Helvetica", 7.5)
        canvas.setFillColor(colors.HexColor(MUTED))
        canvas.drawString(doc_.leftMargin, 7 * mm, f"Repo Change Dashboard · generated {generated}")
        canvas.drawRightString(page[0] - doc_.rightMargin, 7 * mm, f"Page {canvas.getPageNumber()}")
        canvas.restoreState()

    doc.build(story, onFirstPage=furniture, onLaterPages=furniture)
    return buf.getvalue()


# --------------------------------------------------------------------------
# PPTX
# --------------------------------------------------------------------------


def build_pptx(roll: Rollup) -> bytes:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Emu, Inches, Pt

    def rgb(hex_str: str) -> RGBColor:
        return RGBColor.from_string(hex_str.lstrip("#").upper())

    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)
    BLANK = prs.slide_layouts[6]
    SW, SH = prs.slide_width, prs.slide_height

    def add_slide():
        return prs.slides.add_slide(BLANK)

    def textbox(slide, left, top, width, height, text, size, *, bold=False,
                color=INK, align=PP_ALIGN.LEFT):
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        para.alignment = align
        run = para.add_run()
        run.text = str(text)
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
        return box

    def header(slide, title, subtitle=""):
        textbox(slide, Inches(0.6), Inches(0.42), SW - Inches(1.2), Inches(0.6),
                title, 26, bold=True, color=ACCENT_DARK)
        if subtitle:
            textbox(slide, Inches(0.6), Inches(1.02), SW - Inches(1.2), Inches(0.4),
                    subtitle, 12, color=MUTED)
        line = slide.shapes.add_shape(1, Inches(0.6), Inches(1.46), SW - Inches(1.2), Emu(9525))
        line.fill.solid()
        line.fill.fore_color.rgb = rgb(LINE)
        line.line.fill.background()
        line.shadow.inherit = False

    def grid(slide, headers, rows, widths, top=Inches(1.75), row_h=Inches(0.32)):
        """A real PowerPoint table so the numbers stay selectable and editable."""
        shape = slide.shapes.add_table(
            len(rows) + 1, len(headers), Inches(0.6), top,
            SW - Inches(1.2), row_h * (len(rows) + 1)
        )
        tbl = shape.table
        total = sum(widths)
        usable = SW - Inches(1.2)
        for i, w in enumerate(widths):
            tbl.columns[i].width = Emu(int(usable * (w / total)))

        for c, name in enumerate(headers):
            cell = tbl.cell(0, c)
            cell.text = str(name)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(11)
            para.font.bold = True
            para.font.color.rgb = rgb(ACCENT_DARK)
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(ACCENT_SOFT)

        for r, row in enumerate(rows, start=1):
            for c, value in enumerate(row):
                cell = tbl.cell(r, c)
                cell.text = str(value)
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(9.5)
                para.font.color.rgb = rgb(INK)
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb("FFFFFF" if r % 2 else BAND.lstrip("#"))
        return tbl

    crit = roll.criteria
    s = roll.summary
    span = f"{_fmt_date(crit.get('since'))} – {_fmt_date(crit.get('until')) or 'today'}"

    # --- 1. title --------------------------------------------------------
    slide = add_slide()
    band = slide.shapes.add_shape(1, 0, 0, SW, Inches(2.9))
    band.fill.solid()
    band.fill.fore_color.rgb = rgb(ACCENT_SOFT)
    band.line.fill.background()
    band.shadow.inherit = False
    textbox(slide, Inches(0.8), Inches(0.95), SW - Inches(1.6), Inches(0.9),
            "Repository Change Report", 40, bold=True, color=ACCENT_DARK)
    textbox(slide, Inches(0.8), Inches(1.95), SW - Inches(1.6), Inches(0.5),
            f"Commit activity · {span}", 16, color=MUTED)

    rows = _criteria_rows(crit)
    if rows:
        grid(slide, ["Filter", "Value"], [[k, _clip(v, 90)] for k, v in rows],
             [22, 78], top=Inches(3.25))
    textbox(slide, Inches(0.8), SH - Inches(0.7), SW - Inches(1.6), Inches(0.4),
            f"Generated {_fmt_datetime(crit.get('generated_at'))}", 10, color=MUTED)

    # --- 2. KPIs ---------------------------------------------------------
    slide = add_slide()
    header(slide, "Summary", span)
    tiles = [
        ("Commits", s["commits"]), ("Repositories", s["repositories"]),
        ("Services", s["services"]), ("Branches", s["branches"]),
        ("Contributors", s["contributors"]), ("Not on default", s["off_default"]),
        ("Files changed", s["files_changed"]),
    ]
    cols = 4
    box_w = (SW - Inches(1.2) - Inches(0.3) * (cols - 1)) / cols
    for i, (label, value) in enumerate(tiles):
        col, row = i % cols, i // cols
        left = Inches(0.6) + (box_w + Inches(0.3)) * col
        top = Inches(2.0) + (Inches(1.75) + Inches(0.25)) * row
        card = slide.shapes.add_shape(1, left, top, box_w, Inches(1.75))
        card.fill.solid()
        card.fill.fore_color.rgb = rgb("FFFFFF")
        card.line.color.rgb = rgb(LINE)
        card.line.width = Pt(0.75)
        card.shadow.inherit = False
        textbox(slide, left, top + Inches(0.28), box_w, Inches(0.8),
                f"{value:,}", 34, bold=True, color=INK, align=PP_ALIGN.CENTER)
        textbox(slide, left, top + Inches(1.08), box_w, Inches(0.4),
                label, 12, color=MUTED, align=PP_ALIGN.CENTER)

    # --- 3-5. breakdowns -------------------------------------------------
    if roll.by_repo:
        slide = add_slide()
        header(slide, "Activity by repository", span)
        grid(slide, ["Repository", "Commits", "Services touched", "Contributors", "Files"],
             [[_clip(n, 44), v["commits"],
               _join(sorted(x.split(" / ")[-1] for x in v["services"]), 5),
               len(v["authors"]), v["files"]] for n, v in roll.by_repo[:12]],
             [26, 10, 40, 12, 12])

    if roll.by_service:
        slide = add_slide()
        header(slide, "Activity by service / folder", span)
        grid(slide, ["Service (repository / folder)", "Commits", "Contributors", "Files"],
             [[_clip(n, 60), v["commits"], len(v["authors"]), v["files"]]
              for n, v in roll.by_service[:12]],
             [55, 15, 15, 15])

    if roll.by_author:
        slide = add_slide()
        header(slide, "Activity by contributor", span)
        grid(slide, ["Contributor", "Commits", "Repositories", "Services", "Files"],
             [[_clip(n, 40), v["commits"], len(v["repos"]), len(v["services"]), v["files"]]
              for n, v in roll.by_author[:12]],
             [40, 15, 15, 15, 15])

    # --- 6+. commit detail ----------------------------------------------
    detail = roll.commits[:MAX_DETAIL_ROWS]
    per_slide = 11
    pages = [detail[i : i + per_slide] for i in range(0, len(detail), per_slide)]
    for idx, chunk in enumerate(pages, start=1):
        slide = add_slide()
        suffix = f" ({idx}/{len(pages)})" if len(pages) > 1 else ""
        header(slide, f"Commit detail{suffix}", span)
        grid(slide, ["Date", "Repository", "Service", "Branch", "Author", "SHA", "Message"],
             [[_fmt_date(c.get("date")), _clip(c.get("repo"), 24),
               _join(c.get("folders") or [], 2), _join(c.get("branches") or [], 2),
               _clip(c.get("author_name"), 18), (c.get("sha") or "")[:7],
               _clip(c.get("title"), 58)] for c in chunk],
             [9, 15, 14, 12, 12, 7, 31])

    if roll.truncated:
        slide = add_slide()
        header(slide, "Note on scope")
        textbox(
            slide, Inches(0.6), Inches(2.0), SW - Inches(1.2), Inches(2.0),
            f"This deck lists the {MAX_DETAIL_ROWS} most recent commits of "
            f"{len(roll.commits)} matching the filter. The remaining {roll.truncated} "
            "are included in every total and breakdown, but are not listed "
            "individually. Narrow the date range or the filters to list them all.",
            14, color=MUTED,
        )

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
