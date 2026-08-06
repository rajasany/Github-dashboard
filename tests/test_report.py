"""Tests for the date-range window and the report rollup.

Run with:  .venv/bin/python tests/test_report.py
"""

from __future__ import annotations

import sys
from datetime import datetime, timezone
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from app import report  # noqa: E402
from app.main import _resolve_window  # noqa: E402

FAILURES: list[str] = []


def check(name: str, got, want) -> None:
    ok = got == want
    print(f"{'PASS' if ok else 'FAIL'}  {name}")
    if not ok:
        print(f"        got  {got!r}\n        want {want!r}")
        FAILURES.append(name)


def commit(date, repo, folders, branches, author, files=1, on_default=True, sha="abc1234", title="t"):
    return {
        "date": date, "repo": repo, "folders": folders, "branches": branches,
        "author_name": author, "files_changed": files, "on_default": on_default,
        "sha": sha, "title": title, "provider": "github", "url": "#",
    }


COMMITS = [
    commit("2026-08-04T10:00:00Z", "acme/shop", ["backend"], ["main"], "alice", files=3),
    commit("2026-08-04T09:00:00Z", "acme/shop", ["backend", "web"], ["main"], "bob", files=5),
    commit("2026-08-03T09:00:00Z", "acme/shop", ["web"], ["feat/x"], "alice", files=2, on_default=False),
    commit("2026-08-03T08:00:00Z", "gcp/pay", ["backend"], ["master"], "carol", files=7),
    commit("2026-08-01T08:00:00Z", "gcp/pay", ["ledger"], ["master"], "carol", files=1),
]


def test_window() -> None:
    print("=== date-range window ===")

    since, until = _resolve_window(None, "2026-07-01", "2026-08-05")
    check("since starts at midnight UTC", since, datetime(2026, 7, 1, 0, 0, tzinfo=timezone.utc))
    check(
        "until covers the whole named day (not midnight, which would drop it)",
        (until.year, until.month, until.day, until.hour, until.minute),
        (2026, 8, 5, 23, 59),
    )

    since, until = _resolve_window(None, "2026-07-01", None)
    check("open-ended range leaves until unset", until, None)
    check("explicit since is honoured", since.date().isoformat(), "2026-07-01")

    since, until = _resolve_window(30, None, None)
    span = (datetime.now(timezone.utc) - since).days
    check("days-only lookback still works", span, 30)

    # A lookback counted back from an explicit end date, not from today.
    since, until = _resolve_window(10, None, "2026-06-30")
    check("days counts back from `until`, not from now", since.date().isoformat(), "2026-06-20")

    for bad_since, bad_until, label in [
        ("nope", None, "invalid since"),
        (None, "2026-13-45", "invalid until"),
        ("2026-08-01", "2026-07-01", "until before since"),
    ]:
        try:
            _resolve_window(None, bad_since, bad_until)
            check(f"{label} rejected", "no error", "HTTPException")
        except Exception as exc:  # HTTPException
            check(f"{label} rejected", getattr(exc, "status_code", None), 422)


def test_rollup() -> None:
    print("\n=== report rollup ===")
    roll = report.aggregate(COMMITS, {"since": "2026-08-01", "until": None})
    s = roll.summary

    check("commit count", s["commits"], 5)
    check("repository count", s["repositories"], 2)
    # backend appears in both repos and must count twice: (acme/shop, backend)
    # and (gcp/pay, backend) are different services.
    check("services are (repo, folder) pairs", s["services"], 4)
    check("branches are scoped per repo", s["branches"], 3)
    check("contributors", s["contributors"], 3)
    check("not on default", s["off_default"], 1)
    check("files changed is summed", s["files_changed"], 18)

    by_repo = dict(roll.by_repo)
    check("per-repo commit counts", by_repo["acme/shop"]["commits"], 3)
    check("per-repo file counts", by_repo["gcp/pay"]["files"], 8)
    check("per-repo contributor sets", len(by_repo["acme/shop"]["authors"]), 2)

    by_service = dict(roll.by_service)
    check(
        "same-named folders in different repos stay separate",
        (by_service["acme/shop / backend"]["commits"], by_service["gcp/pay / backend"]["commits"]),
        (2, 1),
    )

    by_author = dict(roll.by_author)
    check("author spanning two services", len(by_author["alice"]["services"]), 2)
    check("author commit count", by_author["carol"]["commits"], 2)

    check("days are grouped and newest first", [d for d, _ in roll.by_day], ["2026-08-04", "2026-08-03", "2026-08-01"])
    check("per-day counts", dict(roll.by_day)["2026-08-04"], 2)
    check("nothing truncated for a small set", roll.truncated, 0)

    big = report.aggregate(COMMITS * 120, {})  # 600 commits
    check("truncation is reported, not silent", big.truncated, 600 - report.MAX_DETAIL_ROWS)
    check("totals still cover every commit", big.summary["commits"], 600)


def test_criteria_and_naming() -> None:
    print("\n=== criteria & filenames ===")
    rows = report._criteria_rows(
        {"range": "1 Aug – today", "repository": "acme/shop", "author": None, "search": "fix"}
    )
    check("only set filters are listed", rows, [("Date range", "1 Aug – today"), ("Repository", "acme/shop"), ("Search", "fix")])

    check(
        "filename carries scope and range",
        report.report_filename({"since": "2026-07-01", "until": "2026-08-05", "repository": "acme/shop"}, "pdf"),
        "repo-changes_acme-shop_2026-07-01_to_2026-08-05.pdf",
    )
    check(
        "open-ended range reads as 'now'",
        report.report_filename({"since": "2026-07-01"}, "pptx"),
        "repo-changes_all-repos_2026-07-01_to_now.pptx",
    )
    check("long titles are clipped with an ellipsis", report._clip("x" * 30, 10), "xxxxxxxxx…")
    check("short titles are left alone", report._clip("hello", 10), "hello")
    check("lists over the limit get a +N suffix", report._join(["a", "b", "c", "d"], 2), "a, b +2")
    check("empty lists render as a dash", report._join([]), "—")


def test_documents() -> None:
    print("\n=== document generation ===")
    roll = report.aggregate(COMMITS, {"since": "2026-08-01", "range": "1 Aug – today", "generated_at": "2026-08-06T00:00:00Z"})

    pdf = report.build_pdf(roll)
    check("PDF has a PDF header", pdf[:5], b"%PDF-")
    check("PDF is non-trivial in size", len(pdf) > 3000, True)

    pptx = report.build_pptx(roll)
    check("PPTX is a zip container", pptx[:2], b"PK")

    from io import BytesIO

    from pptx import Presentation

    prs = Presentation(BytesIO(pptx))
    check("deck is 16:9", (round(prs.slide_width / 914400, 2), round(prs.slide_height / 914400, 2)), (13.33, 7.5))
    check("deck has title, summary, 3 breakdowns and detail", len(prs.slides) >= 6, True)

    # Shape order is an implementation detail, so collect every string on the deck.
    all_text = " | ".join(
        shape.text_frame.text
        for slide in prs.slides
        for shape in slide.shapes
        if shape.has_text_frame and shape.text_frame.text.strip()
    )
    for heading in ["Repository Change Report", "Summary", "Activity by repository",
                    "Activity by service / folder", "Activity by contributor", "Commit detail"]:
        check(f"deck contains '{heading}'", heading in all_text, True)
    check("KPI values reach the deck", "18" in all_text and "5" in all_text, True)

    # Every table cell must be non-empty where a value was supplied.
    empties = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_table", False) and shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text == "":
                            empties += 1
    check("no blank cells in any table", empties, 0)


def main() -> int:
    test_window()
    test_rollup()
    test_criteria_and_naming()
    test_documents()
    print(f"\n{'ALL PASS' if not FAILURES else f'{len(FAILURES)} FAILURES: ' + ', '.join(FAILURES)}")
    return 1 if FAILURES else 0


if __name__ == "__main__":
    raise SystemExit(main())
